"""Render a .pptx to PNGs with python-pptx + PIL.

LibreOffice cannot load any .pptx in this container, so this renders the real
generated file instead of a parallel preview. Liberation Sans is metrically
identical to Arial, so text-fit is a trustworthy proxy.
"""
import sys, os, glob, io
from pptx import Presentation
from pptx.util import Emu
from PIL import Image, ImageDraw, ImageFont

def io_bytes(sh): return io.BytesIO(sh.image.blob)

SRC = sys.argv[1]
OUT = sys.argv[2] if len(sys.argv) > 2 else 'render'
SCALE = float(sys.argv[3]) if len(sys.argv) > 3 else 96.0   # px per inch

REG = '/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf'
BLD = '/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf'
_fc = {}
def font(sz, bold):
    k = (int(sz * 2), bold)
    if k not in _fc:
        _fc[k] = ImageFont.truetype(BLD if bold else REG, max(6, int(round(sz))))
    return _fc[k]

def emu_px(v):  return int(round(Emu(v).inches * SCALE))
def rgb(c, d=(0, 0, 0)):
    try:
        if c is None or c.rgb is None: return d
        return tuple(bytes.fromhex(str(c.rgb)))
    except Exception:
        return d

def wrap(draw, text, f, maxw):
    lines = []
    for para in text.split('\n'):
        if not para: lines.append(''); continue
        cur = ''
        for word in para.split(' '):
            t = word if not cur else cur + ' ' + word
            if draw.textlength(t, font=f) <= maxw or not cur:
                cur = t
            else:
                lines.append(cur); cur = word
        lines.append(cur)
    return lines

prs = Presentation(SRC)
SW, SH = emu_px(prs.slide_width), emu_px(prs.slide_height)
for old in glob.glob(f'{OUT}-*.png'): os.remove(old)

for idx, slide in enumerate(prs.slides, 1):
    img = Image.new('RGB', (SW, SH), 'white')
    dr = ImageDraw.Draw(img, 'RGBA')
    # slide background
    try:
        f = slide.background.fill
        if f.type is not None and f.type == 1:
            dr.rectangle([0, 0, SW, SH], fill=rgb(f.fore_color, (255, 255, 255)))
    except Exception:
        pass

    for sh in slide.shapes:
        try:
            x, y = emu_px(sh.left), emu_px(sh.top)
            w, h = emu_px(sh.width), emu_px(sh.height)
        except Exception:
            continue

        if sh.shape_type == 13 or sh.__class__.__name__ == 'Picture':      # image
            try:
                im = Image.open(io_bytes(sh)).convert('RGB')
                iw, ih = im.size
                sc = max(w / iw, h / ih)
                im = im.resize((max(1, int(iw * sc)), max(1, int(ih * sc))), Image.LANCZOS)
                im = im.crop(((im.width - w) // 2, (im.height - h) // 2,
                              (im.width - w) // 2 + w, (im.height - h) // 2 + h))
                img.paste(im, (x, y))
            except Exception as e:
                dr.rectangle([x, y, x + w, y + h], outline=(255, 0, 255), width=2)
            continue

        # filled shape
        try:
            fl = sh.fill
            if fl.type is not None and fl.type == 1:
                col = rgb(fl.fore_color, (200, 200, 200))
                alpha = 255
                el = fl.fore_color._xFill.find(
                    '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                if el is not None:
                    a = el.find('{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                    if a is not None: alpha = int(int(a.get('val')) / 100000 * 255)
                dr.rectangle([x, y, x + w, y + h], fill=col + (alpha,))
        except Exception:
            pass

        if not sh.has_text_frame: continue
        tf = sh.text_frame
        A = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
        runs, chunks = [], []
        for pi, para in enumerate(tf.paragraphs):
            if pi: chunks.append('\n')
            for node in para._p:
                if node.tag == A + 'br': chunks.append('\n')
                elif node.tag == A + 'r':
                    t = node.find(A + 't')
                    chunks.append(t.text or '' if t is not None else '')
            for r in para.runs: runs.append((r.text, r.font))
        if not runs: continue
        para0 = tf.paragraphs[0]
        align = str(para0.alignment) if para0.alignment is not None else 'LEFT (1)'
        sz = 18.0; bold = False; col = (0, 0, 0)
        for _, fo in runs:
            if fo.size: sz = fo.size.pt
            if fo.bold is not None: bold = fo.bold
            break
        fnt = font(sz * SCALE / 72.0, bold)
        text = ''.join(chunks)
        lines = wrap(dr, text, fnt, max(4, w))
        lnspc = 1.22
        sp = para0._p.find(A + 'pPr')
        if sp is not None:
            ls = sp.find(A + 'lnSpc')
            if ls is not None:
                pct = ls.find(A + 'spcPct')
                if pct is not None: lnspc = int(pct.get('val')) / 100000.0
        lh = int(fnt.size * lnspc * 1.02)
        th = lh * len(lines)
        va = str(tf.vertical_anchor)
        ty = y if 'TOP' in va or va == 'None' else (y + (h - th) // 2 if 'MIDDLE' in va else y + h - th)
        overflow = th > h + 2
        for i, ln in enumerate(lines):
            lw = dr.textlength(ln, font=fnt)
            lx = x if 'LEFT' in align else (x + (w - lw) / 2 if 'CENTER' in align else x + w - lw)
            # colour: use first run's colour
            c = (0, 0, 0)
            for _, fo in runs:
                c = rgb(fo.color, (0, 0, 0)); break
            dr.text((lx, ty + i * lh), ln, font=fnt, fill=c)
        if overflow:
            dr.rectangle([x, y, x + w, y + h], outline=(255, 0, 255), width=3)

    img.save(f'{OUT}-{idx:02d}.png')

print('rendered', len(prs.slides._sldIdLst), 'slides at', SW, 'x', SH)
