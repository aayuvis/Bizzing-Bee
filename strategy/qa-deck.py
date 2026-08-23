"""Geometry QA without a renderer.

LibreOffice is broken in this container, so slides cannot be rasterised. These checks
catch the defect classes a visual pass would: out-of-bounds shapes, text that cannot
fit its box, thin margins, and overlapping text frames.

Text fit is ESTIMATED from average glyph width, so it flags candidates rather than
proving overflow. Treated as a warning, not a verdict.
"""
from pptx import Presentation
from pptx.util import Emu
import sys

pr = Presentation('Bizzing-Strategy-Org.pptx')
SW, SH = pr.slide_width, pr.slide_height
E = 914400.0
MARGIN = 0.5

def inches(v): return v / E

issues = []
for i, sl in enumerate(pr.slides, 1):
    frames = []
    for sh in sl.shapes:
        if sh.left is None: continue
        x, y = inches(sh.left), inches(sh.top)
        w, h = inches(sh.width), inches(sh.height)
        # out of slide
        if x < -0.05 or y < -0.05 or x + w > inches(SW) + 0.05 or y + h > inches(SH) + 0.05:
            issues.append(f"S{i}: OFF-SLIDE {sh.shape_type} at ({x:.2f},{y:.2f}) {w:.2f}x{h:.2f}")
        if not sh.has_text_frame: continue
        txt = sh.text_frame.text.strip()
        if not txt: continue
        frames.append((x, y, w, h, txt[:28]))
        # margins
        if x < MARGIN - 0.02 or x + w > inches(SW) - MARGIN + 0.02:
            issues.append(f"S{i}: TIGHT MARGIN x={x:.2f} w={w:.2f} :: {txt[:32]!r}")
        # crude fit estimate
        sizes = [r.font.size.pt for p_ in sh.text_frame.paragraphs for r in p_.runs
                 if r.font.size is not None]
        if sizes:
            pt = max(sizes)
            char_w = pt * 0.50 / 72.0          # avg glyph width in inches
            line_h = pt * 1.34 / 72.0
            per_line = max(1, int(w / char_w))
            lines = 0
            for para in txt.split('\n'):
                lines += max(1, -(-len(para) // per_line))
            need = lines * line_h
            if need > h * 1.12:
                issues.append(f"S{i}: MAY OVERFLOW {need:.2f}in needed vs {h:.2f}in box "
                              f"({pt:.0f}pt) :: {txt[:38]!r}")
    # text-frame overlaps
    for a in range(len(frames)):
        for b in range(a + 1, len(frames)):
            ax, ay, aw, ah, at = frames[a]; bx, by, bw, bh, bt = frames[b]
            ox = min(ax+aw, bx+bw) - max(ax, bx)
            oy = min(ay+ah, by+bh) - max(ay, by)
            if ox > 0.12 and oy > 0.12:
                issues.append(f"S{i}: TEXT OVERLAP {ox:.2f}x{oy:.2f}in :: {at!r} / {bt!r}")

print(f"{len(pr.slides)} slides, {inches(SW):.2f}x{inches(SH):.2f}in")
for t in issues: print(" ", t)
print(f"\n{len(issues)} issue(s)")
